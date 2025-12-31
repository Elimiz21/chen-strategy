#!/usr/bin/env python3
"""
Create Investor Presentation PowerPoint
========================================

Converts the Investor Presentation Deck markdown to a polished PowerPoint file.
Uses python-pptx to create professional slides with consistent styling.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Professional blue theme
COLORS = {
    'primary': RGBColor(0, 82, 147),      # Deep blue
    'secondary': RGBColor(0, 120, 215),    # Bright blue
    'accent': RGBColor(0, 164, 239),       # Light blue
    'success': RGBColor(16, 137, 62),      # Green
    'warning': RGBColor(255, 140, 0),      # Orange
    'text_dark': RGBColor(33, 33, 33),     # Near black
    'text_light': RGBColor(255, 255, 255), # White
    'bg_light': RGBColor(245, 248, 250),   # Light gray
    'highlight': RGBColor(255, 215, 0),    # Gold
}


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with professional styling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add blue header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle if provided
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12.33), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items, highlight_box=None):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            p.text = f"  {item[0]}"
            p.font.size = Pt(item[1])
            p.font.bold = item[2] if len(item) > 2 else False
        else:
            p.text = f"  {item}"
            p.font.size = Pt(20)

        p.font.color.rgb = COLORS['text_dark']
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    # Optional highlight box
    if highlight_box:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_light']
        box.line.color.rgb = COLORS['accent']
        box.line.width = Pt(2)

        box_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(11.93), Inches(0.8))
        tf = box_text.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_box
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_metrics_slide(prs, title, metrics):
    """Add a slide with key metrics in boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Metric boxes
    num_metrics = len(metrics)
    box_width = 3.8
    spacing = (13.33 - (num_metrics * box_width)) / (num_metrics + 1)

    for i, (label, value, subtext) in enumerate(metrics):
        left = spacing + i * (box_width + spacing)

        # Box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(box_width), Inches(2.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['bg_light']
        box.line.color.rgb = COLORS['accent']
        box.line.width = Pt(3)

        # Value
        val_box = slide.shapes.add_textbox(Inches(left), Inches(2.3), Inches(box_width), Inches(1.0))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(left), Inches(3.3), Inches(box_width), Inches(0.6))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_dark']
        p.alignment = PP_ALIGN.CENTER

        # Subtext
        sub_box = slide.shapes.add_textbox(Inches(left), Inches(3.9), Inches(box_width), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['success']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_table_slide(prs, title, headers, rows, highlight_row=None):
    """Add a slide with a data table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = 12.0
    col_width = table_width / num_cols
    row_height = 0.6

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.67), Inches(1.6),
        Inches(table_width), Inches(row_height * num_rows)
    ).table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)

            if highlight_row and i == highlight_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['bg_light']

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['text_dark']
            p.alignment = PP_ALIGN.CENTER

            # Bold first column
            if j == 0:
                p.font.bold = True

    return slide


def add_comparison_slide(prs, title, left_title, right_title, comparisons):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.7)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = COLORS['text_dark']
    left_header.line.fill.background()

    lh_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.6))
    tf = lh_text.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.03), Inches(1.5), Inches(5.8), Inches(0.7)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = COLORS['success']
    right_header.line.fill.background()

    rh_text = slide.shapes.add_textbox(Inches(7.03), Inches(1.55), Inches(5.8), Inches(0.6))
    tf = rh_text.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Comparison rows
    y_pos = 2.4
    for left_val, right_val, label in comparisons:
        # Label
        label_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.33), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_dark']

        # Left value
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.35), Inches(5.8), Inches(0.5))
        tf = left_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_val
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_dark']
        p.alignment = PP_ALIGN.CENTER

        # Right value
        right_box = slide.shapes.add_textbox(Inches(7.03), Inches(y_pos + 0.35), Inches(5.8), Inches(0.5))
        tf = right_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_val
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['success']
        p.alignment = PP_ALIGN.CENTER

        y_pos += 0.95

    return slide


def add_process_slide(prs, title, steps):
    """Add a process/flow slide with connected boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Steps
    num_steps = len(steps)
    box_width = 2.2
    spacing = 0.3
    total_width = num_steps * box_width + (num_steps - 1) * spacing
    start_x = (13.33 - total_width) / 2

    for i, (step_title, step_desc) in enumerate(steps):
        left = start_x + i * (box_width + spacing)

        # Step box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(box_width), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['primary'] if i % 2 == 0 else COLORS['secondary']
        box.line.fill.background()

        # Step number
        num_box = slide.shapes.add_textbox(Inches(left), Inches(2.1), Inches(box_width), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

        # Step title
        title_b = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.6), Inches(box_width - 0.2), Inches(0.8))
        tf = title_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

        # Step description
        desc_b = slide.shapes.add_textbox(Inches(left + 0.1), Inches(3.4), Inches(box_width - 0.2), Inches(1.0))
        tf = desc_b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < num_steps - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(left + box_width + 0.05), Inches(3.1), Inches(0.2), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['accent']
            arrow.line.fill.background()

    return slide


def add_validation_slide(prs, title, tests):
    """Add a validation checklist slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Test items
    y_pos = 1.6
    for test_name, result, status in tests:
        # Check box
        check = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(0.5), Inches(0.5)
        )
        check.fill.solid()
        check.fill.fore_color.rgb = COLORS['success']
        check.line.fill.background()

        check_text = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.05), Inches(0.5), Inches(0.4))
        tf = check_text.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

        # Test name
        name_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.05), Inches(4.5), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = test_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_dark']

        # Result
        result_box = slide.shapes.add_textbox(Inches(5.7), Inches(y_pos + 0.05), Inches(4.5), Inches(0.5))
        tf = result_box.text_frame
        p = tf.paragraphs[0]
        p.text = result
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_dark']

        # Status badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(y_pos + 0.05), Inches(2.0), Inches(0.45)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLORS['success']
        badge.line.fill.background()

        badge_text = slide.shapes.add_textbox(Inches(10.5), Inches(y_pos + 0.1), Inches(2.0), Inches(0.4))
        tf = badge_text.text_frame
        p = tf.paragraphs[0]
        p.text = status
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER

        y_pos += 0.75

    return slide


def add_cta_slide(prs, title, subtitle, points):
    """Add a call-to-action closing slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.7))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    # Points
    points_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.5), Inches(9.33), Inches(2.5))
    tf = points_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_light']
        p.space_before = Pt(12)

    # Contact box
    contact = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.67), Inches(5.8), Inches(6.0), Inches(0.8)
    )
    contact.fill.solid()
    contact.fill.fore_color.rgb = COLORS['highlight']
    contact.line.fill.background()

    contact_text = slide.shapes.add_textbox(Inches(3.67), Inches(5.95), Inches(6.0), Inches(0.6))
    tf = contact_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Contact us to discuss investment opportunity"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the full investor presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    add_title_slide(prs, "ADAPTIVE REGIME-AWARE TRADING SYSTEM", "QQQ (Nasdaq-100 ETF) | Investor Presentation | December 2025")

    # Slide 2: Investment Thesis
    add_content_slide(prs, "Investment Thesis: Systematic Alpha in QQQ", [
        ("KEY INSIGHT", 24, True),
        "Different trading strategies perform best in different market conditions.",
        "By detecting market regimes in real-time and dynamically allocating",
        "across validated strategies, we generate exceptional risk-adjusted returns.",
        "",
        ("OUR APPROACH", 24, True),
        "DETECT market regime (100 micro-regimes, 4 dimensions)",
        "SELECT optimal strategies (7 validated, beat academic baselines)",
        "ALLOCATE capital dynamically ('Tilt Not Switch' philosophy)",
    ])

    # Slide 3: Performance Headlines
    add_metrics_slide(prs, "Exceptional Risk-Adjusted Returns", [
        ("SHARPE RATIO", "8.78", "175x vs Buy-and-Hold"),
        ("ANNUAL RETURN", "114.6%", "14.7x vs Buy-and-Hold"),
        ("MAX DRAWDOWN", "-2.7%", "13x lower risk"),
    ])

    # Slide 4: Problem Statement
    add_content_slide(prs, "The Problem: Traditional Approaches Fall Short", [
        ("BUY-AND-HOLD", 20, True),
        "  35% drawdowns destroy capital and investor psychology",
        "  Years to recover from major losses",
        "",
        ("SIMPLE TECHNICAL ANALYSIS", 20, True),
        "  Fails after transaction costs",
        "  No statistical edge vs random entries",
        "",
        ("SINGLE STRATEGY SYSTEMS", 20, True),
        "  Works in some regimes, fails catastrophically in others",
        "  No adaptation to changing market conditions",
        "",
        ("TRADITIONAL REGIME DETECTION", 20, True),
        "  Only 2 regimes (bull/bear) - too coarse",
        "  1,000+ day duration - not actionable",
    ], highlight_box="OUR SOLUTION: 4-dimensional micro-regime detection with 100 states and 1.8-day average duration")

    # Slide 5: Innovation - Micro-Regime Detection
    add_content_slide(prs, "Our Innovation: 867x More Granular Regime Detection", [
        ("FOUR DIMENSIONS OF MARKET STATE", 24, True),
        "",
        "1. TREND (5 levels): Strong Bull → Bull → Neutral → Bear → Strong Bear",
        "   Captures direction and strength of market movement",
        "",
        "2. VOLATILITY (4 levels): Low → Normal → High → Crisis",
        "   Identifies risk environment and optimal position sizing",
        "",
        "3. MOMENTUM (3 levels): Accelerating → Steady → Decelerating",
        "   Detects trend strength changes before they occur",
        "",
        "4. MEAN-REVERSION (3 levels): Overbought → Neutral → Oversold",
        "   Measures extension from fair value",
    ], highlight_box="RESULT: 100 distinct micro-regimes | 1.8 days average duration | 95.8% day-over-day stability")

    # Slide 6: Comparison
    add_comparison_slide(prs, "From Strategic to Tactical Signals",
        "TRADITIONAL", "OUR APPROACH",
        [
            ("1 dimension (trend only)", "4 dimensions", "Dimensionality"),
            ("2 regimes (bull/bear)", "100 micro-regimes", "Granularity"),
            ("1,000+ days", "1.8 days", "Average Duration"),
            ("Strategic (wait years)", "Tactical (daily signals)", "Actionability"),
            ("—", "867x IMPROVEMENT", "Enhancement Factor"),
        ]
    )

    # Slide 7: Strategy Selection Process
    add_process_slide(prs, "Rigorous Strategy Validation Process", [
        ("Build Universe", "21 TA strategies\n6 academic baselines"),
        ("Academic Benchmark", "Must beat\nTrendEnsemble\n(Sharpe 3.88)"),
        ("Statistical Testing", "Walk-forward\nMonte Carlo\nMultiple testing"),
        ("Regime Stability", "Must perform\nacross ALL regimes"),
        ("Final Selection", "7 strategies\nvalidated"),
    ])

    # Slide 8: Winning Strategies
    add_table_slide(prs, "5 Strategies Beat Best Academic Baseline",
        ["Rank", "Strategy", "Sharpe", "vs Baseline", "Status"],
        [
            ["1", "BBSqueeze", "10.61", "+6.72", "SELECTED"],
            ["2", "DonchianBreakout", "8.18", "+4.30", "SELECTED"],
            ["3", "KeltnerBreakout", "5.55", "+1.67", "SELECTED"],
            ["4", "Ichimoku", "5.00", "+1.12", "SELECTED"],
            ["5", "ParabolicSAR", "4.56", "+0.68", "SELECTED"],
            ["—", "TrendEnsemble (Baseline)", "3.88", "0.00", "Benchmark"],
            ["—", "QQQ Buy-and-Hold", "0.05", "-3.83", "—"],
        ]
    )

    # Slide 9: Portfolio Construction
    add_table_slide(prs, "7-Strategy Regime-Aware Portfolio",
        ["Strategy", "Weight", "Role"],
        [
            ["BBSqueeze", "25%", "Primary alpha source"],
            ["DonchianBreakout", "25%", "Trend following, excellent in BEAR"],
            ["KeltnerBreakout", "15%", "Volatility breakouts"],
            ["TrendEnsemble", "10%", "Academic diversification"],
            ["RORO", "10%", "Defensive protection"],
            ["Ichimoku", "10%", "Trend confirmation"],
            ["ParabolicSAR", "5%", "Trend following"],
        ]
    )

    # Slide 10: Robust Across Periods
    add_table_slide(prs, "Consistent Performance Across All Market Periods",
        ["Period", "Years", "Sharpe", "QQQ Context"],
        [
            ["Dot-com Crash", "2000-2005", "7.57", "QQQ fell -78%"],
            ["Financial Crisis", "2006-2009", "7.20", "Lehman collapse, -50% crash"],
            ["Post-Crisis Bull", "2010-2015", "7.86", "Recovery and QE"],
            ["Late Cycle", "2016-2019", "6.92", "Extended bull run"],
            ["COVID & Beyond", "2020-2024", "7.77", "Pandemic volatility"],
        ]
    )

    # Slide 11: Zero Overfitting
    add_content_slide(prs, "Zero Overfitting Evidence: Out-of-Sample > In-Sample", [
        ("THE 'OVERFITTING PROBLEM' IN QUANT FINANCE", 22, True),
        "Most trading systems fail live because they're overfit to historical data.",
        "They 'memorize' past patterns that don't repeat.",
        "",
        ("OUR RESULTS: NEGATIVE SHARPE DECAY", 22, True),
        "",
        "In-Sample Sharpe:      8.62",
        "Out-of-Sample Sharpe:  8.76",
        "Sharpe Decay:          -1.7% (NEGATIVE = GOOD)",
        "",
        ("INTERPRETATION", 22, True),
        "The model captures REAL market dynamics, not random noise.",
        "This negative decay is statistically RARE and indicates genuine edge.",
    ])

    # Slide 12: Independent Validation
    add_validation_slide(prs, "6-Point Independent Validation (All Passed)", [
        ("1. REPLICATION", "Exact match from clean environment", "PASSED"),
        ("2. SUBPERIOD STABILITY", "All 5 periods Sharpe > 6.9, CV = 4.8%", "PASSED"),
        ("3. PARAMETER SENSITIVITY", "All configurations Sharpe > 7.0", "PASSED"),
        ("4. COST SENSITIVITY", "Profitable at 5x costs (Sharpe 7.64)", "PASSED"),
        ("5. DRAWDOWN CONSTRAINTS", "All scenarios < 10% DD", "PASSED"),
        ("6. MODEL RISK", "No critical risks identified", "PASSED"),
    ])

    # Slide 13: Risk Management
    add_content_slide(prs, "Multi-Layer Risk Management Framework", [
        ("LAYER 1: STRATEGY LEVEL", 20, True),
        "  Individual strategy kill criteria",
        "  Strategy suspended if underperforms by >50%",
        "",
        ("LAYER 2: PORTFOLIO LEVEL", 20, True),
        "  20% Maximum Drawdown (hard limit)",
        "  Dynamic leverage reduction:",
        "    DD < 10%: Full leverage",
        "    DD 10-15%: Reduce to 0.75x",
        "    DD 15-18%: Reduce to 0.50x",
        "    DD > 20%: Force liquidation + 5-day cooldown",
        "",
        ("LAYER 3: SYSTEM LEVEL", 20, True),
        "  Real-time monitoring dashboard",
        "  Automated alerts at 3%, 5%, 10% DD",
    ], highlight_box="HISTORICAL MAX DD: -2.67% (well under 20% limit)")

    # Slide 14: Cost Sensitivity
    add_table_slide(prs, "Profitable Even at 5x Transaction Costs",
        ["Cost Multiplier", "Sharpe Ratio", "Annual Return", "Status"],
        [
            ["1.0x (Baseline)", "8.95", "170.1%", "PROFITABLE"],
            ["2.0x Costs", "8.66", "164.4%", "PROFITABLE"],
            ["3.0x Costs", "8.39", "158.7%", "PROFITABLE"],
            ["5.0x Costs", "7.64", "141.0%", "PROFITABLE"],
        ]
    )

    # Slide 15: Development Process
    add_table_slide(prs, "Phase-Gated Development Rigor",
        ["Phase", "Description", "Status"],
        [
            ["0", "Charter & Success Definition", "PASSED"],
            ["1", "Literature & Design Space Map", "PASSED"],
            ["2", "Data Foundation & Research Stack", "PASSED"],
            ["3", "Expert Library & Cost Model", "PASSED"],
            ["4", "Regime Detection & Calibration", "PASSED"],
            ["5", "Meta-Allocation Engine", "PASSED"],
            ["6", "INDEPENDENT VALIDATION", "PASSED"],
            ["7", "Paper Trading (30 days)", "PENDING"],
            ["8", "Live Pilot ($500K)", "PENDING"],
        ]
    )

    # Slide 16: Expert Review
    add_content_slide(prs, "Academic-Grade Methodology Validation", [
        ("EXPERT PERSPECTIVES INCORPORATED", 22, True),
        "",
        "Dr. Marcos Lopez de Prado (Advances in Financial ML)",
        "  Multiple testing correction, purged cross-validation",
        "",
        "Dr. David Aronson (Evidence-Based Technical Analysis)",
        "  Data mining bias prevention, statistical significance",
        "",
        "Dr. Andrew Lo (MIT - Adaptive Markets Hypothesis)",
        "  Regime-dependent strategy validation",
        "",
        "Dr. Nassim Taleb (Black Swan author)",
        "  Tail risk and stress testing methodology",
    ], highlight_box="KEY OUTCOME: Expert panel recommended 100 micro-regime framework (867x improvement)")

    # Slide 17: Key Differentiators
    add_content_slide(prs, "What Makes This Unique", [
        ("1. 867x MORE GRANULAR REGIME DETECTION", 18, True),
        "   Tactical signals (1.8-day regimes) vs strategic (1,000+ day)",
        "",
        ("2. ACADEMIC BASELINE BENCHMARKING", 18, True),
        "   Proves GENUINE alpha, not comparison to naive buy-and-hold",
        "",
        ("3. ZERO OVERFITTING EVIDENCE", 18, True),
        "   OOS > IS performance (rare) - captures real market dynamics",
        "",
        ("4. SIX-POINT VALIDATION", 18, True),
        "   Replication, subperiod, parameter, cost, DD, model risk - all passed",
        "",
        ("5. PHASE-GATED GOVERNANCE", 18, True),
        "   No premature deployment, explicit gate criteria required",
        "",
        ("6. FULL TRANSPARENCY", 18, True),
        "   All code, data, decisions documented & reproducible",
    ])

    # Slide 18: Investment Parameters
    add_table_slide(prs, "Investment Parameters",
        ["Parameter", "Value"],
        [
            ["Initial Capital", "$500,000"],
            ["Asset", "QQQ (Nasdaq-100 ETF)"],
            ["Strategy", "Long/Short QQQ + Cash, 7-strategy allocation"],
            ["Maximum Leverage", "2.0x (dynamically managed)"],
            ["Maximum Drawdown Limit", "25% ($125,000 max loss)"],
            ["Target Performance", "Sharpe > 1.0 (achieved: 8.78)"],
            ["Signal Frequency", "Daily (after market close)"],
            ["Minimum Hold Period", "1 day (no intraday trading)"],
        ]
    )

    # Slide 19: Summary
    add_content_slide(prs, "Investment Summary", [
        ("EXCEPTIONAL RETURNS", 22, True),
        "  Sharpe 8.78 | 114.6% Annual Return | -2.67% Max DD",
        "",
        ("RIGOROUS VALIDATION", 22, True),
        "  6/6 validation tests passed | Zero overfitting evidence",
        "  Stable across ALL market periods (2000-2024)",
        "",
        ("ROBUST RISK MANAGEMENT", 22, True),
        "  Multi-layer defense system | 20% hard DD limit (achieved 2.67%)",
        "  Dynamic leverage adjustment based on drawdown",
        "",
        ("INNOVATIVE APPROACH", 22, True),
        "  867x more granular regime detection",
        "  Beats sophisticated academic baselines",
        "  Full transparency and reproducibility",
    ])

    # Slide 20: Call to Action
    add_cta_slide(prs, "Ready for Deployment",
        "Phase 7: Paper Trading | Phase 8: Live Pilot ($500K)",
        [
            "Paper Trading: 30 days to validate live execution",
            "Live Pilot: $500K capital, 90-day track record",
            "Maximum Risk: $125,000 (25% hard limit)",
            "Full transparency and daily reporting",
        ]
    )

    return prs


def main():
    """Generate the PowerPoint presentation."""
    print("Creating Investor Presentation PowerPoint...")

    prs = create_presentation()

    # Save to docs folder
    output_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
        "docs",
        "Investor_Presentation.pptx"
    )

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
